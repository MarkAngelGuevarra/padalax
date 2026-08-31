import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation in 16:9 Widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Theme Colors
BG_DARK = RGBColor(11, 15, 25)         # Deep Slate / Dark Navy
CARD_BG = RGBColor(18, 26, 43)         # Card background
CARD_BORDER = RGBColor(30, 41, 59)     # Border
CYAN = RGBColor(6, 182, 212)           # Primary Neon Cyan
BLUE = RGBColor(59, 130, 246)          # Accent Blue
EMERALD = RGBColor(16, 185, 129)       # Success Emerald
AMBER = RGBColor(245, 158, 11)         # Warning / Gold
ROSE = RGBColor(244, 63, 94)           # Danger / Contrast Rose
WHITE = RGBColor(255, 255, 255)        # Pure White
SLATE_200 = RGBColor(226, 232, 240)    # Light gray text
SLATE_400 = RGBColor(148, 163, 184)    # Muted gray text

def set_slide_background(slide, color):
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()

def add_header(slide, title_text, category_text="PADALAX PROTOCOL • STELLAR & SOROBAN"):
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = CYAN
    p_cat.font.name = "Arial"

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    p_title.font.name = "Arial"

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.2)
    return card

# 1. TITLE
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1, BG_DARK)
add_card(slide1, 0.8, 1.0, 11.733, 5.5, CARD_BG, CYAN)

badge = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.5), Inches(4.8), Inches(0.45))
badge.fill.solid()
badge.fill.fore_color.rgb = RGBColor(6, 45, 60)
badge.line.color.rgb = CYAN
badge.line.width = Pt(1)
p_b = badge.text_frame.paragraphs[0]
p_b.text = "STELLAR RISEIN • LEVEL 5 BLUE BELT SUBMISSION"
p_b.font.size = Pt(11)
p_b.font.bold = True
p_b.font.color.rgb = CYAN
p_b.font.name = "Arial"
p_b.alignment = PP_ALIGN.CENTER

title_box = slide1.shapes.add_textbox(Inches(1.3), Inches(2.1), Inches(10.7), Inches(1.8))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PadalaX Protocol"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Arial"

p2 = tf.add_paragraph()
p2.text = "Decentralized Cross-Border Remittance & Voucher Escrow on Stellar"
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = CYAN
p2.font.name = "Arial"

desc_box = slide1.shapes.add_textbox(Inches(1.3), Inches(4.0), Inches(10.7), Inches(1.2))
tf_d = desc_box.text_frame
tf_d.word_wrap = True
p_d = tf_d.paragraphs[0]
p_d.text = (
    "Slashing international remittance fees from traditional 8% down to < $0.001 with sub-5-second settlement. "
    "Senders lock cryptographic vouchers in Soroban smart contracts; unbanked families cash out directly to GCash, Maya, and local bank accounts."
)
p_d.font.size = Pt(14)
p_d.font.color.rgb = SLATE_200
p_d.font.name = "Arial"

foot_box = slide1.shapes.add_textbox(Inches(1.3), Inches(5.4), Inches(10.7), Inches(0.6))
tf_f = foot_box.text_frame
p_f = tf_f.paragraphs[0]
p_f.text = "Live App: padalax.vercel.app  •  GitHub: github.com/MarkAngelGuevarra/padalax  •  Author: Mark Guevarra"
p_f.font.size = Pt(12)
p_f.font.bold = True
p_f.font.color.rgb = EMERALD
p_f.font.name = "Arial"

# 2. PROBLEM
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2, BG_DARK)
add_header(slide2, "The Problem: The $2.2 Billion Annual Predatory Fee Drain", "MARKET PAIN POINTS")

problems = [
    ("Predatory 5%-8.5% Fees", "Traditional money transfer operators (Western Union, MoneyGram, banks) extract massive cuts from hard-earned OFW wages, draining families of essential capital.", ROSE),
    ("Slow 24-72 Hour Clearance", "Legacy wire transfers and bank corridors take days to settle. Families experiencing medical, tuition, or food emergencies are stranded waiting for confirmations.", AMBER),
    ("44% Unbanked Divide", "Over 44% of Filipino adults lack traditional bank accounts or Web3 crypto wallets, creating severe onboarding friction for non-technical recipients in provinces.", CYAN),
]

for i, (p_title, p_desc, p_color) in enumerate(problems):
    left = 0.8 + i * 3.97
    add_card(slide2, left, 1.7, 3.75, 4.8, CARD_BG, p_color)
    pill = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.3), Inches(2.0), Inches(3.15), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(30, 20, 35) if p_color == ROSE else (RGBColor(35, 30, 20) if p_color == AMBER else RGBColor(15, 35, 45))
    pill.line.color.rgb = p_color
    p_p = pill.text_frame.paragraphs[0]
    p_p.text = f"PAIN POINT #{i+1}"
    p_p.font.size = Pt(10)
    p_p.font.bold = True
    p_p.font.color.rgb = p_color
    p_p.alignment = PP_ALIGN.CENTER
    
    t_box = slide2.shapes.add_textbox(Inches(left + 0.3), Inches(2.6), Inches(3.15), Inches(0.8))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_p = t_tf.paragraphs[0]
    t_p.text = p_title
    t_p.font.size = Pt(17)
    t_p.font.bold = True
    t_p.font.color.rgb = WHITE
    
    d_box = slide2.shapes.add_textbox(Inches(left + 0.3), Inches(3.5), Inches(3.15), Inches(2.7))
    d_tf = d_box.text_frame
    d_tf.word_wrap = True
    d_p = d_tf.paragraphs[0]
    d_p.text = p_desc
    d_p.font.size = Pt(13)
    d_p.font.color.rgb = SLATE_200

# 3. SOLUTION
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3, BG_DARK)
add_header(slide3, "The Solution: Cryptographic Vouchers & Instant Fiat Settlement", "VALUE PROPOSITION")

solutions = [
    ("Sub-Cent Remittance Fees", "Slashes transaction fees from $15-$35 down to < 0.00001 XLM (< $0.001) per transfer. OFWs keep 99.99% of their money.", EMERALD),
    ("Sub-5-Second Settlement", "Powered by the Stellar Consensus Protocol (SCP). Funds arrive in Manila within seconds, not days.", BLUE),
    ("Secret PIN Voucher Escrow", "Senders lock funds into Soroban with SHA-256 hashlocks. Recipients claim via WhatsApp/SMS PIN directly to GCash/Maya with 100% time-locked autonomous refund.", CYAN),
]

for i, (s_title, s_desc, s_color) in enumerate(solutions):
    left = 0.8 + i * 3.97
    add_card(slide3, left, 1.7, 3.75, 4.8, CARD_BG, s_color)
    pill = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.3), Inches(2.0), Inches(3.15), Inches(0.4))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(10, 35, 30) if s_color == EMERALD else (RGBColor(15, 25, 45) if s_color == BLUE else RGBColor(15, 35, 45))
    pill.line.color.rgb = s_color
    p_p = pill.text_frame.paragraphs[0]
    p_p.text = f"CORE ADVANTAGE #{i+1}"
    p_p.font.size = Pt(10)
    p_p.font.bold = True
    p_p.font.color.rgb = s_color
    p_p.alignment = PP_ALIGN.CENTER
    
    t_box = slide3.shapes.add_textbox(Inches(left + 0.3), Inches(2.6), Inches(3.15), Inches(0.8))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    t_p = t_tf.paragraphs[0]
    t_p.text = s_title
    t_p.font.size = Pt(17)
    t_p.font.bold = True
    t_p.font.color.rgb = WHITE
    
    d_box = slide3.shapes.add_textbox(Inches(left + 0.3), Inches(3.5), Inches(3.15), Inches(2.7))
    d_tf = d_box.text_frame
    d_tf.word_wrap = True
    d_p = d_tf.paragraphs[0]
    d_p.text = s_desc
    d_p.font.size = Pt(13)
    d_p.font.color.rgb = SLATE_200

# 4. MARKET
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4, BG_DARK)
add_header(slide4, "Market Opportunity: $37.2 Billion Inflow & 10M+ OFWs", "GLOBAL MARKET SIZE")

add_card(slide4, 0.8, 1.7, 5.6, 4.8, CARD_BG, CYAN)
stat_box = slide4.shapes.add_textbox(Inches(1.1), Inches(2.0), Inches(5.0), Inches(4.2))
stf = stat_box.text_frame
stf.word_wrap = True
p = stf.paragraphs[0]
p.text = "Philippines Remittance Market"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN

p = stf.add_paragraph()
p.text = "$37.2 BILLION"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = WHITE

p = stf.add_paragraph()
p.text = "Annual Inflows (Top 4 Globally)\n"
p.font.size = Pt(13)
p.font.color.rgb = EMERALD

p = stf.add_paragraph()
p.text = (
    "• 10M+ Overseas Filipino Workers worldwide\n"
    "• $2.2B Lost Annually to intermediary fees\n"
    "• 76M+ Active GCash & Maya accounts ready for instant SEP-24 mobile off-ramping."
)
p.font.size = Pt(12)
p.font.color.rgb = SLATE_200

corridors = [
    ("UAE & Saudi Arabia", "42% of Global OFW Remittances", "Primary corridor for nurses, construction engineers, and domestic professionals."),
    ("Singapore & Hong Kong", "28% of Regional Inflows", "High-frequency monthly allowance senders with smartphone literacy."),
    ("USA & Canada", "20% of High-Ticket Transfers", "College tuition, real estate amortization, and medical support."),
    ("UK & Europe", "10% of European Corridor", "Growing healthcare and maritime seafarer population."),
]

for i, (c_title, c_sub, c_desc) in enumerate(corridors):
    col = i % 2
    row = i // 2
    c_left = 6.65 + col * 2.92
    c_top = 1.7 + row * 2.45
    add_card(slide4, c_left, c_top, 2.82, 2.35, CARD_BG, BLUE)
    c_box = slide4.shapes.add_textbox(Inches(c_left + 0.15), Inches(c_top + 0.15), Inches(2.52), Inches(2.0))
    c_tf = c_box.text_frame
    c_tf.word_wrap = True
    cp = c_tf.paragraphs[0]
    cp.text = c_title
    cp.font.size = Pt(12)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp2 = c_tf.add_paragraph()
    cp2.text = c_sub
    cp2.font.size = Pt(10)
    cp2.font.bold = True
    cp2.font.color.rgb = CYAN
    cp3 = c_tf.add_paragraph()
    cp3.text = c_desc
    cp3.font.size = Pt(9.5)
    cp3.font.color.rgb = SLATE_400

# 5. ARCHITECTURE
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5, BG_DARK)
add_header(slide5, "System Architecture: Soroban Token Escrow & SHA-256 Hashlocks", "TECHNICAL ARCHITECTURE")

arch_steps = [
    ("1. OFW Sender", "Connects Freighter/Albedo\nSelects Amount (XLM/USDC)\nCalculates Live PHP Rate", BLUE),
    ("2. Soroban Escrow", "Contract ID: CATUXAJ7...\nReal token::Client Transfer\nDataKey::Remittance(id)\nextend_ttl Persistence", CYAN),
    ("3. Secret PIN", "SHA-256 Preimage Hash\nOne-time 6-digit Claim PIN\n1-Click WhatsApp/Telegram", AMBER),
    ("4. Recipient Payout", "Zero-Gas Web Claim\nInstant GCash / Maya Payout\nOr Autonomous Sender Refund", EMERALD),
]

for i, (a_title, a_desc, a_color) in enumerate(arch_steps):
    left = 0.8 + i * 2.98
    add_card(slide5, left, 1.7, 2.8, 3.2, CARD_BG, a_color)
    box = slide5.shapes.add_textbox(Inches(left + 0.15), Inches(1.9), Inches(2.5), Inches(2.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = a_title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = a_color
    p2 = tf.add_paragraph()
    p2.text = a_desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = SLATE_200

add_card(slide5, 0.8, 5.15, 11.733, 1.4, CARD_BG, CARD_BORDER)
sum_box = slide5.shapes.add_textbox(Inches(1.1), Inches(5.25), Inches(11.1), Inches(1.2))
stf = sum_box.text_frame
stf.word_wrap = True
p = stf.paragraphs[0]
p.text = "Security & Invariants Verified on Testnet"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = CYAN
p2 = stf.add_paragraph()
p2.text = (
    "• 6/6 Rust Unit Tests Passing (test_create_and_claim, test_refund_after_expiration, unauthorized attack panics)\n"
    "• Autonomous Time-Locked Refund: Senders retain 100% cryptographic authority to claw back unclaimed expired vouchers.\n"
    "• Storage Footprint & TTL Optimization: Storage bumps configured to protect persistent ledger state."
)
p2.font.size = Pt(10.5)
p2.font.color.rgb = SLATE_200

# 6. TRACTION
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6, BG_DARK)
add_header(slide6, "Traction & Pilot Validation: 50+ Verified On-Chain Users", "LEVEL 5 VALIDATION METRICS")

metrics = [
    ("50+", "Verified Pilot Users", "Onboarded across 25 global OFW cities (Dubai, Riyadh, Singapore, Tokyo, London, Toronto)", CYAN),
    ("100%", "Transaction Success", "50/50 on-chain remittance escrows confirmed on Stellar Testnet Horizon", EMERALD),
    ("4.85 / 5.0", "Avg User Satisfaction", "Collected via structured Google Form feedback surveys across all 50 participants", AMBER),
    ("PHP 34,500+", "Fees Saved in Pilot", "Compared to standard 6.5% traditional remittance center charges", BLUE),
]

for i, (m_val, m_title, m_desc, m_color) in enumerate(metrics):
    left = 0.8 + i * 2.98
    add_card(slide6, left, 1.7, 2.8, 4.8, CARD_BG, m_color)
    box = slide6.shapes.add_textbox(Inches(left + 0.15), Inches(2.0), Inches(2.5), Inches(4.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = m_val
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = m_color
    p2 = tf.add_paragraph()
    p2.text = m_title
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p3 = tf.add_paragraph()
    p3.text = f"\n{m_desc}"
    p3.font.size = Pt(11)
    p3.font.color.rgb = SLATE_200

# 7. FEEDBACK EVOLUTION
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7, BG_DARK)
add_header(slide7, "Product Iteration: Built Directly From User Feedback", "FEEDBACK-DRIVEN EVOLUTION")

feedbacks = [
    ("Feedback 1: Multi-Family Allowance", "Users needed to send money to multiple family members at once.", "Multi-Recipient Batch Remittance: Create up to 5 vouchers in 1 flow (Commit a8627a6).", CYAN),
    ("Feedback 2: Currency Conversion Friction", "OFWs in UAE/Saudi struggled with manual AED/SAR to PHP math.", "Live Multi-Fiat FX Ticker: Real-time rates & mini-estimator for 10 currencies (Commit f4e1d17).", BLUE),
    ("Feedback 3: Messaging App Sharing", "Copying links was cumbersome for non-technical parents.", "1-Click WhatsApp & Telegram Share: Direct voucher link dispatch & PDF receipt (Commit 7767804).", EMERALD),
    ("Feedback 4: Error Telemetry", "Required monitoring for wallet disconnections and RPC failures.", "Vercel Analytics & Error Boundary: Embedded root telemetry (Commit 2e165cb).", AMBER),
]

for i, (f_title, f_pain, f_sol, f_color) in enumerate(feedbacks):
    top = 1.7 + i * 1.25
    add_card(slide7, 0.8, top, 11.733, 1.15, CARD_BG, f_color)
    box = slide7.shapes.add_textbox(Inches(1.0), Inches(top + 0.1), Inches(11.3), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{f_title}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = f_color
    p2 = tf.add_paragraph()
    p2.text = f"User Request: {f_pain}  -->  Implemented: {f_sol}"
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = SLATE_200

# 8. ROADMAP
slide8 = prs.slides.add_slide(blank_layout)
set_slide_background(slide8, BG_DARK)
add_header(slide8, "Roadmap: Path to Mainnet & Ecosystem Integration", "FUTURE ROADMAP")

roadmap = [
    ("Levels 1-3 (White, Yellow, Orange)", "ACCEPTED", "Stellar SDK, Soroban Escrow Contract, 6 Unit Tests, Testnet Deployment.", EMERALD),
    ("Level 4 (Green Belt)", "COMPLETED", "Production Web3 PWA on Vercel, Vercel Analytics, 12 Pilot Users Onboarded.", EMERALD),
    ("Level 5 (Blue Belt)", "SUBMITTED", "50+ Real Verified Testnet Users, Batch Remittance, Multi-Fiat FX Ticker, Pitch Deck.", CYAN),
    ("Level 6 (Black Belt)", "NEXT PHASE", "Stellar Mainnet Deployment & Gasless Relayer Fee Sponsorship.", BLUE),
    ("Level 7 (Master Track)", "VISION", "Live SEP-24 Regulated Off-Ramping with Philippine Banking Anchors (GCash & Maya).", AMBER),
]

for i, (r_title, r_badge, r_desc, r_color) in enumerate(roadmap):
    top = 1.7 + i * 1.02
    add_card(slide8, 0.8, top, 11.733, 0.92, CARD_BG, r_color)
    box = slide8.shapes.add_textbox(Inches(1.0), Inches(top + 0.08), Inches(11.3), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{r_title}  •  [{r_badge}]"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = r_color
    p2 = tf.add_paragraph()
    p2.text = r_desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = SLATE_200

downloads_dir = os.path.join(os.path.expanduser("~"), "Downloads")
pptx_path = os.path.join(downloads_dir, "PadalaX_Pitch_Deck.pptx")
prs.save(pptx_path)
print("PowerPoint presentation successfully saved at: " + pptx_path)

scratch_path = os.path.join(os.path.expanduser("~"), ".gemini", "antigravity", "scratch", "padalax", "PadalaX_Pitch_Deck.pptx")
prs.save(scratch_path)
print("Copy saved in repo at: " + scratch_path)
